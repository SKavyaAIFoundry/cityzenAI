from pptx import Presentation                   # MIT License
from pptx.enum.shapes import MSO_SHAPE_TYPE


# method for extracting content from PowerPoint files

def analysePptxFile(fileName):
    # https://stackoverflow.com/questions/39418620/extracting-text-from-multiple-powerpoint-files-using-python
    # https://python-pptx.readthedocs.io/en/latest/api/enum/MsoShapeType.html
    word_count = 0
    paragraph_count = 0
    bullet_count = 0
    image_count = 0

    prs = Presentation(fileName)
    powerpoint_text = []

    # extract text from all elements of all slides and add to list
    # and count the number of images
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                powerpoint_text.append(shape.text + " ") # " " is break words at the end of the current shape and start of the next shape
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                image_count += 1

    print(powerpoint_text)

    document_text = ''.join(powerpoint_text)

    # create temp text file and populate with document text
    temp_file = './data/powerpoint_text.txt'
    with open(temp_file, 'w') as text_file:
        text_file.write(document_text)

    # read text from temp text file for analysis
    document = open(temp_file)
    document_text = document.read()

    paragraph_count = document_text.count('\n')
    paragraph_count += bullet_count
    
    document_words = document_text.split()
    word_count = len(document_words)

    return document_text, word_count, paragraph_count, bullet_count, image_count
